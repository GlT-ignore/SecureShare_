import os
import bcrypt
import json
import mysql.connector
from flask import Flask, render_template, request, redirect, url_for, session, flash, send_file, after_this_request, Response
from Crypto.Cipher import AES
from Crypto.Util.Padding import pad, unpad
from Crypto.Hash import SHA256
from Crypto.PublicKey import DSA
from Crypto.Signature import DSS
from flask import abort
import logging
from datetime import datetime, timedelta
import re
from functools import wraps
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import io
import csv
from io import StringIO
import xlsxwriter
import traceback

def check_required_files():
    required_files = [
        (DSA_PRIVATE_KEY_PATH, 'DSA private key'),
        (DSA_PUBLIC_KEY_PATH, 'DSA public key'),
        (AES_KEY_PATH, 'AES key')
    ]
    
    missing_files = []
    for file_path, file_desc in required_files:
        if not os.path.exists(file_path):
            missing_files.append(file_desc)
    
    if missing_files:
        print("Error: Missing required files:")
        for file_desc in missing_files:
            print(f"- {file_desc}")
        print("\nPlease run the following scripts:")
        print("python signatures.py")
        print("python generate_aes_key.py")
        exit(1)

# First define the app
app = Flask(__name__)
app.secret_key = 'your_secret_key'

# Get the absolute path to the project directory
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Set up configuration BEFORE checking for the upload folder
app.config['UPLOAD_FOLDER'] = os.path.join(BASE_DIR, 'uploads')
AES_KEY_PATH = os.path.join(BASE_DIR, 'aes_key.bin')
DSA_PRIVATE_KEY_PATH = os.path.join(BASE_DIR, 'dsa_private_key.pem')
DSA_PUBLIC_KEY_PATH = os.path.join(BASE_DIR, 'dsa_public_key.pem')

# Now check and create the upload folder if it doesn't exist
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

db_config = {
    # 'host': 'mysql-18825a65-siddharth11-8758.e.aivencloud.com',
    # 'user': 'avnadmin',
    # 'password': 'AVNS_yAkAzRO3mqqBdknAU9t',
    # 'database': 'is_project',
    # 'port': '21868'

    'host' : "127.0.0.1",
'user': "root",
'password' : "",
'database' : "is_project",
'port' : 3306
}


def init_db():
    conn = mysql.connector.connect(**db_config)
    return conn


def close_db(conn):
    conn.close()


def read_aes_key():
    with open(AES_KEY_PATH, 'rb') as file:
        key = file.read(16)
    return key


def aes_encrypt(key, file_data):
    cipher = AES.new(key, AES.MODE_CBC)
    iv = cipher.iv
    ciphertext = cipher.encrypt(pad(file_data, AES.block_size))
    return iv, ciphertext


def aes_decrypt(key, iv, ciphertext):
    cipher = AES.new(key, AES.MODE_CBC, iv)
    decrypted = unpad(cipher.decrypt(ciphertext), AES.block_size)
    return decrypted


def hash_keyword(keyword):
    hasher = SHA256.new()
    hasher.update(keyword.encode('utf-8'))
    return hasher.hexdigest()

def sign(data):
    with open(DSA_PRIVATE_KEY_PATH, 'rb') as key_file:
        private_key = DSA.import_key(key_file.read())
    data_hash = SHA256.new(data)
    signer = DSS.new(private_key, 'fips-186-3')
    signature = signer.sign(data_hash)
    return signature


def verify(data, signature):
    with open(DSA_PUBLIC_KEY_PATH, 'rb') as key_file:
        public_key = DSA.import_key(key_file.read())
    data_hash = SHA256.new(data)
    verifier = DSS.new(public_key, 'fips-186-3')
    try:
        verifier.verify(data_hash, signature)
        return True
    except (ValueError, TypeError):
        return False


def index_file(file_id, keywords):
    conn = init_db()
    cursor = conn.cursor()
    for keyword in keywords:
        hashed_keyword = hash_keyword(keyword)
        try:
            insert_query = "INSERT INTO FileIndex (keyword, file_id) VALUES (%s, %s)"
            cursor.execute(insert_query, (hashed_keyword, file_id))
        except mysql.connector.Error as err:
            print(f"Error inserting keyword '{keyword}': {err}")
    conn.commit()
    cursor.close()
    close_db(conn)


def upload_file(file_path):
    try:
        # Log the start of upload process
        log_action("Debug", "Starting upload process")
        
        # Get file details before reading
        filesize = os.path.getsize(file_path)
        filename = os.path.basename(file_path)
        
        # Read and encrypt the file
        with open(file_path, 'rb') as file:
            file_data = file.read()
            
        key = read_aes_key()
        iv, ciphertext = aes_encrypt(key, file_data)
        signature = sign(ciphertext)
        
        # Get usernames
        target_username = request.form.get('username')
        sender_username = session.get('username')
        
        log_action("Debug", f"Sending file from {sender_username} to {target_username}")
        
        # Store in database
        conn = init_db()
        cursor = conn.cursor()
        
        try:
            # Check if recipient exists
            cursor.execute("SELECT username FROM Users WHERE username = %s", (target_username,))
            if not cursor.fetchone():
                log_action("Upload Error", f"Recipient {target_username} does not exist")
                flash(f"User '{target_username}' does not exist", "error")
                cursor.close()
                close_db(conn)
                return 0
                
            # Insert file
            insert_query = """
                INSERT INTO Files (username, sender, filename, iv, encrypted_file, 
                                 signature, filesize, upload_time) 
                VALUES (%s, %s, %s, %s, %s, %s, %s, NOW())
            """
            cursor.execute(insert_query, (
                target_username,
                sender_username,
                filename,
                iv,
                ciphertext,
                signature,
                filesize
            ))
            
            file_id = cursor.lastrowid
            conn.commit()
            
            log_action("Debug", f"File uploaded successfully. ID: {file_id}")
            flash(f"File successfully sent to {target_username}", "success")
            return 1
            
        except mysql.connector.Error as err:
            log_action("Database Error", f"Failed to insert file: {str(err)}")
            flash("Error uploading file", "error")
            return 0
            
        finally:
            cursor.close()
            close_db(conn)
            
    except Exception as e:
        log_action("Upload Error", f"Error in upload_file: {str(e)}\n{traceback.format_exc()}")
        flash("Error processing file", "error")
        return 0


def extract_text_content(file_path):
    file_ext = os.path.splitext(file_path)[1].lower()
    content = ""
    
    try:
        if file_ext == '.pdf':
            try:
                # Log the PDF processing attempt
                log_action("PDF Processing", f"Attempting to process PDF file: {file_path}")
                
                # Check if file exists and is readable
                if not os.path.exists(file_path):
                    log_action("PDF Error", f"PDF file does not exist: {file_path}")
                    return ""
                    
                # Try to read the PDF file
                with open(file_path, 'rb') as pdf_file:
                    try:
                        # Create PDF reader
                        pdf_reader = PdfReader(pdf_file)
                        page_count = len(pdf_reader.pages)
                        log_action("PDF Processing", f"Successfully opened PDF with {page_count} pages")
                        
                        # Extract text from each page
                        for page_num in range(page_count):
                            try:
                                page = pdf_reader.pages[page_num]
                                page_text = page.extract_text()
                                if page_text:
                                    content += page_text + "\n"
                                log_action("PDF Processing", f"Successfully extracted text from page {page_num + 1}")
                            except Exception as page_error:
                                log_action("PDF Error", f"Error on page {page_num + 1}: {str(page_error)}")
                                continue
                                
                    except Exception as reader_error:
                        log_action("PDF Error", f"Error creating PDF reader: {str(reader_error)}")
                        return ""
                        
            except Exception as file_error:
                log_action("PDF Error", f"Error opening PDF file: {str(file_error)}")
                return ""
                
            # Log the final content length
            log_action("PDF Processing", f"Extracted {len(content)} characters of text")
            
        elif file_ext == '.docx':
            doc = Document(file_path)
            for para in doc.paragraphs:
                content += para.text + "\n"
                
        elif file_ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
                        
        else:  # For text files
            encodings = ['utf-8', 'latin-1', 'ascii']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
                    
        return content
        
    except Exception as e:
        log_action("Content Extraction Error", f"Failed to extract content from {file_path}: {str(e)}\n{traceback.format_exc()}")
        return ""
        
    return content


def receive_files(username):
    log_action("Debug", f"Fetching received files for user: {username}")
    conn = init_db()
    cursor = conn.cursor()
    select_query = """
        SELECT file_id, sender, filename, iv, encrypted_file, signature, upload_time, filesize 
        FROM Files 
        WHERE username = %s 
        ORDER BY upload_time DESC
    """
    try:
        cursor.execute(select_query, (username,))
        results = cursor.fetchall()
        log_action("Debug", f"Found {len(results)} received files")

        verified_files = []
        for file in results:
            file_id, sender, filename, iv, encrypted_file, signature, upload_time, filesize = file
            verified = verify(encrypted_file, signature)
            adjusted_time = upload_time - timedelta(hours=0, minutes=0)
            verified_files.append((file_id, sender, filename, iv, encrypted_file, signature, verified, adjusted_time, filesize))
            log_action("Debug", f"Processed file: {filename} from {sender}")

        return verified_files

    except Exception as e:
        log_action("Error", f"Error fetching received files: {str(e)}\n{traceback.format_exc()}")
        return []
    finally:
        cursor.close()
        close_db(conn)


def search_encrypted_files(keyword):
    log_action("Search", f"Admin {session.get('username')} searched for keyword: {keyword}")
    conn = init_db()
    cursor = conn.cursor()

    hashed_keyword = hash_keyword(keyword)
    cursor.execute("SELECT file_id FROM FileIndex WHERE keyword = %s", (hashed_keyword,))
    file_ids = cursor.fetchall()

    cursor.close()
    close_db(conn)

    return [file_id[0] for file_id in file_ids]


def check_credentials(username, password):
    conn = init_db()
    cursor = conn.cursor()
    cursor.execute("SELECT password, role FROM Users WHERE username = %s", (username,))
    result = cursor.fetchone()
    close_db(conn)

    if result:
        hashed_password, role = result
        if bcrypt.checkpw(password.encode('utf-8'), hashed_password.encode('utf-8')):
            return role
    return None


@app.route('/')
def home():
    # If user is already logged in, redirect to appropriate dashboard
    if 'username' in session:
        if session.get('role') == 'admin':
            return redirect(url_for('admin_page'))
        return redirect(url_for('user_page'))
    return render_template('login.html')


@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']

        role = check_credentials(username, password)

        if role:
            log_action("Login", f"User '{username}' logged in successfully with role '{role}'")
            session['username'] = username
            session['role'] = role
            if role == 'admin':
                return redirect(url_for('admin_page'))
            else:
                return redirect(url_for('user_page'))
        else:
            log_action("Failed Login", f"Failed login attempt for username '{username}'")
            flash("Invalid credentials, please try again.", "error")  # Flash the error message
            return redirect(url_for('home'))
    return render_template('login.html')


@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')

        # Validate inputs
        if not username or not password or not confirm_password:
            flash('All fields are required', 'error')
            return render_template('signup.html')

        if password != confirm_password:
            flash('Passwords do not match', 'error')
            return render_template('signup.html')

        try:
            conn = init_db()
            cursor = conn.cursor()

            # Check if username already exists
            cursor.execute("SELECT username FROM Users WHERE username = %s", (username,))
            if cursor.fetchone():
                flash('Username already exists', 'error')
                return render_template('signup.html')

            # Create new user
            hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
            cursor.execute(
                "INSERT INTO Users (username, password, role) VALUES (%s, %s, %s)",
                (username, hashed_password, 'user')
            )
            conn.commit()

            log_action("User Creation", f"New user account created: {username}")
            flash('Account created successfully! Please login.', 'success')
            return redirect(url_for('login'))

        except mysql.connector.Error as err:
            log_action("Failed Registration", f"Failed to create user {username}: {str(err)}")
            flash(f'Error creating account: {str(err)}', 'error')
            return render_template('signup.html')

        finally:
            if 'cursor' in locals():
                cursor.close()
            if 'conn' in locals() and conn.is_connected():
                conn.close()

    # GET request - show signup form
    return render_template('signup.html')


@app.route('/user', methods=['GET', 'POST'])
def user_page():
    message = None
    if request.method == 'POST':
        try:
            # Debug log - Check if file is in request
            log_action("Debug", "POST request received in user_page")
            
            if 'file' not in request.files:
                log_action("Upload Error", "No file part in request")
                message = {"type": "error", "text": "No file selected"}
                return render_template('user.html', message=message)
                
            file = request.files['file']
            target_username = request.form.get('username')
            
            # Debug log - Check file and username
            log_action("Debug", f"File name: {file.filename}, Target username: {target_username}")
            
            if file.filename == '':
                log_action("Upload Error", "No file selected")
                message = {"type": "error", "text": "No file selected"}
                return render_template('user.html', message=message)

            try:
                # Ensure upload directory exists
                if not os.path.exists(app.config['UPLOAD_FOLDER']):
                    os.makedirs(app.config['UPLOAD_FOLDER'])
                
                # Save file with a unique name
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                safe_filename = f"{timestamp}_{file.filename}"
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)
                
                log_action("Debug", f"Saving file to: {file_path}")
                file.save(file_path)
                
                # Process the file
                result = upload_file(file_path)
                log_action("Debug", f"Upload result: {result}")
                
                # Clean up
                if os.path.exists(file_path):
                    os.remove(file_path)
                    log_action("Upload Process", "Temporary file cleaned up")

                if result == 1:
                    message = {"type": "success", "text": "File sent successfully!"}
                else:
                    message = {"type": "error", "text": "Failed to send file. Please check logs for details."}

            except Exception as e:
                log_action("Upload Error", f"Error in file upload process: {str(e)}\n{traceback.format_exc()}")
                if 'file_path' in locals() and os.path.exists(file_path):
                    os.remove(file_path)
                message = {"type": "error", "text": f"Error uploading file: {str(e)}"}

        except Exception as e:
            log_action("Upload Error", f"Unexpected error in user_page: {str(e)}\n{traceback.format_exc()}")
            message = {"type": "error", "text": f"Unexpected error: {str(e)}"}

    # Get files for display
    username = session.get('username')
    log_action("Debug", f"Getting files for username: {username}")
    
    files = receive_files(username)
    log_action("Debug", f"Received files count: {len(files)}")
    
    sent_files = get_sent_files(username)
    log_action("Debug", f"Sent files count: {len(sent_files)}")
    
    return render_template('user.html', message=message, files=files, sent_files=sent_files)


@app.route('/received_files')
def received_files():
    username = session.get('username')
    files = receive_files(username)
    return render_template('received_files.html', files=files)


@app.route('/download_file/<int:file_id>')
def download_file(file_id):
    try:
        conn = init_db()
        cursor = conn.cursor()
        cursor.execute("SELECT iv, encrypted_file, filename FROM Files WHERE file_id = %s", (file_id,))
        result = cursor.fetchone()
        close_db(conn)

        if result:
            iv, encrypted_file, filename = result
            key = read_aes_key()
            decrypted_data = aes_decrypt(key, iv, encrypted_file)

            # Create temp directory if it doesn't exist
            temp_dir = os.path.join(BASE_DIR, 'temp')
            if not os.path.exists(temp_dir):
                os.makedirs(temp_dir)

            # Use absolute path for temp file
            temp_path = os.path.join(temp_dir, f'decrypted_file_{file_id}{os.path.splitext(filename)[1]}')
            
            try:
                with open(temp_path, 'wb') as f:
                    f.write(decrypted_data)

                log_action("File Download", f"File ID: {file_id}, Filename: {filename}, Downloaded by: {session.get('username')}")
                
                @after_this_request
                def remove_file(response):
                    try:
                        if os.path.exists(temp_path):
                            os.remove(temp_path)
                    except Exception as e:
                        log_action("Cleanup Error", f"Error removing temp file {temp_path}: {str(e)}")
                    return response

                return send_file(
                    temp_path,
                    as_attachment=True,
                    download_name=filename,
                    max_age=0
                )
            except Exception as e:
                log_action("Download Error", f"Error processing file {filename}: {str(e)}")
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                return "Error processing file", 500

        return "File not found.", 404

    except Exception as e:
        log_action("Download Error", f"Error downloading file ID {file_id}: {str(e)}")
        return "Error downloading file", 500


def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('role') == 'admin':
            abort(403)  # Forbidden access
        return f(*args, **kwargs)
    return decorated_function


@app.route('/admin', methods=['GET', 'POST'])
@admin_required
def admin_page():
    conn = init_db()
    cursor = conn.cursor()
    
    # Update the Users query to fetch specific fields
    cursor.execute("SELECT username, password, role FROM Users")
    users = cursor.fetchall()
    
    # Update the Files query to fetch specific fields
    cursor.execute("""
        SELECT file_id, filename, username, sender, upload_time, filesize 
        FROM Files 
        ORDER BY upload_time DESC
    """)
    files = cursor.fetchall()
    
    # Update the FileIndex query
    cursor.execute("SELECT keyword, file_id FROM FileIndex")
    file_index = cursor.fetchall()
    
    # Calculate total storage
    cursor.execute("SELECT COALESCE(SUM(filesize), 0) FROM Files")
    total_storage = cursor.fetchone()[0]
    
    cursor.close()
    close_db(conn)
    
    if request.method == 'POST':
        keyword = request.form['keyword']
        file_ids = search_encrypted_files(keyword)
        return render_template('admin.html', 
                            keyword=keyword, 
                            file_ids=file_ids, 
                            users=users, 
                            files=files, 
                            file_index=file_index, 
                            total_storage=total_storage)
    
    return render_template('admin.html', 
                        users=users, 
                        files=files, 
                        file_index=file_index, 
                        total_storage=total_storage)


@app.route('/delete_user', methods=['POST'])
@admin_required
def delete_user():
    username = request.form['username']
    
    try:
        conn = init_db()
        cursor = conn.cursor()

        # First delete user's files from FileIndex
        cursor.execute("""
            DELETE fi FROM FileIndex fi
            INNER JOIN Files f ON fi.file_id = f.file_id
            WHERE f.username = %s OR f.sender = %s
        """, (username, username))

        # Delete user's files
        cursor.execute("DELETE FROM Files WHERE username = %s OR sender = %s", (username, username))
        
        # Delete the user
        cursor.execute("DELETE FROM Users WHERE username = %s AND role != 'admin'", (username,))
        
        conn.commit()
        log_action("User Deletion", f"User {username} deleted by admin {session.get('username')}")
        flash(f"User {username} has been deleted.", "success")
        
    except mysql.connector.Error as err:
        log_action("Failed User Deletion", f"Failed to delete user {username}: {str(err)}")
        flash(f"Error deleting user: {err}", "error")
    finally:
        cursor.close()
        close_db(conn)

    return redirect(url_for('admin_page'))


@app.route('/delete_file', methods=['POST'])
@admin_required
def delete_file():
    file_id = request.form['file_id']
    conn = init_db()
    cursor = conn.cursor()

    try:
        # First delete from FileIndex table
        cursor.execute("DELETE FROM FileIndex WHERE file_id = %s", (file_id,))
        
        # Then delete from Files table
        cursor.execute("DELETE FROM Files WHERE file_id = %s", (file_id,))
        
        conn.commit()
        log_action("File Deletion", f"File ID {file_id} deleted by admin {session.get('username')}")
        flash(f"File ID {file_id} has been deleted.", "success")
    except mysql.connector.Error as err:
        log_action("Failed File Deletion", f"Failed to delete file ID {file_id}: {str(err)}")
        flash(f"Error deleting file: {err}", "error")
    finally:
        cursor.close()
        close_db(conn)

    return redirect(url_for('admin_page'))


@app.route('/logout')
def logout():
    username = session.get('username')
    log_action("Logout", f"User '{username}' logged out")
    session.pop('username', None)
    session.pop('role', None)
    return redirect(url_for('home'))


def log_action(action, details):
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_entry = f"[{timestamp}] {action}: {details}\n"
    with open('logs.txt', 'a') as log_file:
        log_file.write(log_entry)


def get_sent_files(username):
    log_action("Debug", f"Fetching sent files for user: {username}")
    conn = init_db()
    cursor = conn.cursor()
    select_query = """
        SELECT file_id, username, filename, iv, encrypted_file, signature, upload_time, filesize 
        FROM Files 
        WHERE sender = %s 
        ORDER BY upload_time DESC
    """
    try:
        cursor.execute(select_query, (username,))
        results = cursor.fetchall()
        log_action("Debug", f"Found {len(results)} sent files")

        verified_files = []
        for file in results:
            file_id, recipient, filename, iv, encrypted_file, signature, upload_time, filesize = file
            verified = verify(encrypted_file, signature)
            adjusted_time = upload_time - timedelta(hours=0, minutes=0)
            verified_files.append((file_id, recipient, filename, iv, encrypted_file, signature, verified, adjusted_time, filesize))
            log_action("Debug", f"Processed file: {filename} to {recipient}")

        return verified_files

    except Exception as e:
        log_action("Error", f"Error fetching sent files: {str(e)}\n{traceback.format_exc()}")
        return []
    finally:
        cursor.close()
        close_db(conn)


@app.errorhandler(413)
def request_entity_too_large(error):
    flash('File too large. Maximum size is 16MB.', 'error')
    return redirect(url_for('user_page')), 413


@app.route('/export_users_csv')
@admin_required
def export_users_csv():
    conn = init_db()
    cursor = conn.cursor()
    cursor.execute("SELECT username, role FROM Users")
    users = cursor.fetchall()
    cursor.close()
    close_db(conn)

    # Create CSV in memory
    si = StringIO()
    cw = csv.writer(si)
    cw.writerow(['Username', 'Role', 'Status'])  # Headers
    for user in users:
        cw.writerow([user[0], user[1], 'Active'])
    
    output = si.getvalue()
    si.close()

    return Response(
        output,
        mimetype="text/csv",
        headers={"Content-Disposition": f"attachment;filename=users_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"}
    )

@app.route('/export_files_csv')
@admin_required
def export_files_csv():
    conn = init_db()
    cursor = conn.cursor()
    cursor.execute("""
        SELECT filename, username, sender, upload_time, filesize 
        FROM Files 
        ORDER BY upload_time DESC
    """)
    files = cursor.fetchall()
    cursor.close()
    close_db(conn)

    # Create CSV in memory
    si = StringIO()
    cw = csv.writer(si)
    cw.writerow(['File Name', 'Owner', 'Sender', 'Upload Date', 'Size (KB)'])  # Headers
    for file in files:
        cw.writerow([
            file[0], 
            file[1], 
            file[2], 
            file[3].strftime('%Y-%m-%d %H:%M'),
            round(file[4] / 1024, 1)
        ])
    
    output = si.getvalue()
    si.close()

    return Response(
        output,
        mimetype="text/csv",
        headers={"Content-Disposition": f"attachment;filename=files_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"}
    )

@app.route('/export_all_csv')
@admin_required
def export_all_csv():
    conn = init_db()
    cursor = conn.cursor()
    
    # Get users data
    cursor.execute("SELECT username, role FROM Users")
    users = cursor.fetchall()
    
    # Get files data
    cursor.execute("""
        SELECT filename, username, sender, upload_time, filesize 
        FROM Files 
        ORDER BY upload_time DESC
    """)
    files = cursor.fetchall()
    
    cursor.close()
    close_db(conn)
    
    # Create CSV in memory
    si = StringIO()
    cw = csv.writer(si)
    
    # Write system overview
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    total_storage = sum(file[4] for file in files) if files else 0
    
    cw.writerow(['SYSTEM OVERVIEW'])
    cw.writerow(['Report Generated:', timestamp])
    cw.writerow(['Total Users:', len(users)])
    cw.writerow(['Total Files:', len(files)])
    cw.writerow(['Total Storage:', f"{(total_storage / 1024 / 1024):.2f} MB"])
    cw.writerow([])  # Empty row for spacing
    
    # Write users section
    cw.writerow(['USER MANAGEMENT'])
    cw.writerow(['Username', 'Role', 'Status'])
    for user in users:
        cw.writerow([user[0], user[1], 'Active'])
    cw.writerow([])  # Empty row for spacing
    
    # Write files section
    cw.writerow(['FILE MANAGEMENT'])
    cw.writerow(['File Name', 'Owner', 'Sender', 'Upload Date', 'Size (KB)'])
    for file in files:
        cw.writerow([
            file[0],
            file[1],
            file[2],
            file[3].strftime('%Y-%m-%d %H:%M'),
            round(file[4] / 1024, 1)
        ])
    
    output = si.getvalue()
    si.close()
    
    return Response(
        output,
        mimetype="text/csv",
        headers={"Content-Disposition": f"attachment;filename=system_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"}
    )


if __name__ == "__main__":
    app.run(debug=True)
